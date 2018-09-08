from pptx import Presentation
import os
import re

print("pptquiz.py\n\nPer creare un file powerpoint con domande e risposte che appaiono nella slide successiva a quella della domanda, crea un file di testo con la domanda (o il titolo) nella prima riga, con le risposte nella riga successiva tra parentesi quadre, puoi mettere solo una parola per volta, quindi, se la risposta è fatta di 3 parole, metti tre parole tra parentesi quadre. Puoi anche mettere del testo normale e le parole tra le parentesi quadre come parole omesse.\n\nEsempio:\nQual è la capitale d'Italia?\n[Roma]\n\nApparità la domanda nella prima slide e la risposta nella seconda...")


for files in os.listdir():
    if files.endswith(".txt"):
        print(files)
SORGENTE = input("Quale file vuoi caricare? >")
if not SORGENTE.endswith(".txt"):
    SORGENTE = SORGENTE + ".txt"

# SORGENTE = "testo2.txt"
DESTINAZIONE = SORGENTE[:-4] + ".pptx"

prs = Presentation()


def slide(number=0):
    newslide = 0
    title_slide_layout = prs.slide_layouts[number]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = file.readline()
    x = file.readline()
    print(x)
    if "[" in x:
        ptrn = "\[([A-Za-z0-9_\s]+)\]"
        y = x
        words = re.findall(ptrn, x)
        for word in words:
            x = x.replace(word, "_" * len(word))
            newslide = 1
    subtitle.text = x
    if newslide:
        title_slide_layout = prs.slide_layouts[number]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Soluzione"
        subtitle.text = y #+ "\nParole omesse:\n[" + "-".join(words) + "]"
    file.readline()


with open(SORGENTE, 'r', encoding="utf-8") as file:
    file_lenght = len(file.readlines())
    file.seek(0)
    slide(0)
    for i in range(int(file_lenght / 3)):
        slide(0)

try:
    ppt = DESTINAZIONE[:4] + "_quiz" + ".pptx"
    prs.save(ppt)
    import os
    os.system(ppt)
except PermissionError:
    print("Chiudo il file aperto")
    os.system("TASKKILL /IM POWERPNT.EXE")
    os.system(ppt)
