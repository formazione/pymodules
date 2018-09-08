from pptx import Presentation
import os
import seefiles

# pymodules
seefiles.seefiles()

SORGENTE = input("Quale file vuoi caricare? >")
if not SORGENTE.endswith(".txt"):
    SORGENTE = SORGENTE + ".txt"

# SORGENTE = "testo2.txt"
DESTINAZIONE = SORGENTE[:-4] + ".pptx"

prs = Presentation()


def slide(number=0):
    title_slide_layout = prs.slide_layouts[number]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = file.readline()
    file.readline()
    subtitle.text = file.readline()
    file.readline()


with open(SORGENTE, 'r', encoding="utf-8") as file:
    file_lenght = len(file.readlines())
    file.seek(0)
    slide(0)
    for i in range(int(file_lenght / 4)):
        slide(1)

prs.save(DESTINAZIONE)

os.system(DESTINAZIONE)
